"""Generate a 10-slide e& AI Stack presentation matching the on-page brand style."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


BRAND = {
    "red":        RGBColor(0xE0, 0x08, 0x00),
    "red_dark":   RGBColor(0xA8, 0x20, 0x1A),
    "black":      RGBColor(0x11, 0x11, 0x11),
    "white":      RGBColor(0xFF, 0xFF, 0xFF),
    "grey":       RGBColor(0x88, 0x88, 0x88),
    "grey_dark":  RGBColor(0x44, 0x44, 0x44),
    "grey_mid":   RGBColor(0x66, 0x66, 0x66),
    "grey_soft":  RGBColor(0xAA, 0xAA, 0xAA),
    "border":     RGBColor(0xE6, 0xE6, 0xE6),
    "lightgrey":  RGBColor(0xF7, 0xF7, 0xF7),
    "panel_bg":   RGBColor(0xFA, 0xFA, 0xFA),
    "purple":     RGBColor(0x7A, 0x52, 0xF4),
    "blue":       RGBColor(0x18, 0x5F, 0xA5),
    "green":      RGBColor(0x00, 0x4B, 0x2E),
    "rose_bg":    RGBColor(0xFF, 0xF5, 0xF5),
    "forge_a":    RGBColor(0x0A, 0x0A, 0x0A),
    "forge_b":    RGBColor(0x2A, 0x08, 0x08),
}

FONT = "Arial"  # Suisse International fallback


STACK_LAYERS = [
    {
        "n": "6",
        "name": "Customers & Distribution",
        "tag": "Where AI meets revenue",
        "color": BRAND["red"],
        "summary": ("The marketplace, sales force, partner channel, and e& invoice that put "
                    "AI solutions into reach across e&'s 240M+ direct billing relationships "
                    "globally — including 320K+ SMBs, enterprises and government clients in "
                    "core markets."),
        "owner": "e& (100%)",
        "components": ["SMB marketplace", "e& sales force", "e& invoice integration",
                       "Enterprise direct", "Government tenders", "Partner channel",
                       "Multi-OpCo expansion"],
        "why": ("e&'s unfair advantage: 240M+ direct billing relationships globally, telco "
                "trust and brand reach, plus tens of thousands of physical antenna sites "
                "across its markets; in core markets, 320K+ SMBs on monthly billing and "
                "zero-CAC add-ons to the existing e& invoice."),
    },
    {
        "n": "5",
        "name": "Agentic Applications",
        "tag": "What customers actually buy",
        "color": BRAND["red_dark"],
        "summary": ("End-to-end AI agents that own a business function — Customer, Sales, "
                    "Comms, Finance, Ops, People — plus enterprise solutions and sovereign "
                    "deployments. Future roadmap extends to physical-world edge: humanoid "
                    "robots, autonomous drones, connected vehicles, IoT devices."),
        "owner": "e& — agent IP under progressive build-then-transfer",
        "components": ["6 SMB agents", "Tier 1 Adapted Agents", "Tier 2 Custom AI",
                       "Tier 3 Sovereign / On-Prem", "Contact Centre AI",
                       "Document Intelligence", "Approval Orchestration",
                       "Edge: humanoid robots · drones · vehicles (future)"],
        "why": ("Agents are the product. Everything below this layer is the platform that "
                "makes them possible. As e&'s physical network expands, the same agentic "
                "layer will extend into robots, drones, and connected devices."),
    },
    {
        "n": "4",
        "name": "Forge — AI Control Plane + Development",
        "tag": "The durable platform IP",
        "color": BRAND["green"],
        "summary": ("The unified software layer that runs the AI business (control plane) "
                    "and builds new AI fast (SDD). Forge connects compute, models, agents, "
                    "customers, and billing into one operating system."),
        "owner": "AIdeology · licensed to e&",
        "is_forge": True,
        "pillars": [
            {
                "title": "Operations · Control Plane",
                "sub": "How Forge runs the AI business",
                "color": BRAND["green"],
                "items": [
                    ("Agent orchestration", "Multi-agent workflows, shared context, failover, human-in-loop gates"),
                    ("LLM gateway & model router", "One API across Claude, GPT, Llama, sovereign models"),
                    ("GPU & compute orchestration", "Tenant quotas, cost allocation, sovereign placement"),
                    ("Billing & monetization", "Subscriptions, usage metering, revenue share reporting"),
                    ("Governance & compliance", "RBAC, audit, PII masking, trust tiers SMB → Government"),
                    ("Observability", "Traces, cost-per-agent, SLA dashboards, anomaly alerts"),
                ],
            },
            {
                "title": "Development · Spec-Driven Build",
                "sub": "How Forge builds new AI fast",
                "color": BRAND["red"],
                "items": [
                    ("SDD methodology", "Spec-first: signed implementation spec before code"),
                    ("Solution blueprints", "Reusable templates for each agent type"),
                    ("Prompt libraries", "Tested, versioned prompt patterns"),
                    ("Guardrail packs", "Pre-built safety, compliance and quality controls"),
                    ("Connector accelerators", "Salesforce, SAP, MOHRE, e& billing, etc."),
                    ("Acceptance test framework", "Every spec ships with named tests"),
                ],
            },
        ],
    },
    {
        "n": "3",
        "name": "Model Layer",
        "tag": "The AI intelligence",
        "color": BRAND["purple"],
        "summary": ("The models themselves — commercial APIs, open-source, and sovereign "
                    "Arabic-tuned models — served on high-throughput inference and "
                    "fine-tuning infrastructure."),
        "owner": "Mixed: third-party APIs + e&-hosted sovereign models",
        "components": ["GPT-4o (Azure OpenAI)", "Claude (Anthropic)", "Llama · Falcon · Mistral",
                       "Sovereign Arabic-tuned models", "vLLM · TGI · NVIDIA Triton",
                       "DeepSpeed · FSDP · LoRA · QLoRA", "Model registry & versioning",
                       "STT / TTS voice pipelines"],
        "why": ("Models change every six months. The model layer must be swappable. "
                "Forge routes between providers without touching agent code."),
    },
    {
        "n": "2",
        "name": "GPU Orchestration — Open Innovation",
        "tag": "Compute on demand",
        "color": BRAND["blue"],
        "summary": ("e&'s GPU orchestrator. Schedules, isolates, and serves GPU workloads "
                    "— inference, training, fine-tuning — across e&'s sovereign compute "
                    "estate. Orchestration partner: Open Innovation."),
        "owner": "e&",
        "components": ["Kubernetes / Slurm job scheduling", "Multi-tenant GPU isolation",
                       "Inference · training · fine-tuning", "Auto-scaling & load balancing",
                       "Hardware abstraction (DGX, Dell, HPE, Supermicro)",
                       "Capacity management & tenant quotas"],
        "why": ("The GPU orchestrator turns raw hardware into a consumable, schedulable, "
                "billable service. Forge consumes capacity from it and maps every job back "
                "to a customer, agent and invoice line."),
    },
    {
        "n": "1",
        "name": "Infrastructure",
        "tag": "The physical foundation",
        "color": BRAND["black"],
        "summary": ("The hardware, data centres, network, and devices that everything runs "
                    "on — including e&'s edge and physical device footprint."),
        "owner": "e& (100%)",
        "components": ["e& Data Centres · UAE × 3, Morocco, Hungary, PPF",
                       "NVIDIA DGX · Dell · HPE · Supermicro · H100 / B200",
                       "G42 / Core42 · Azure · AWS · OCI",
                       "VAST Data · DDN · Pure FlashBlade storage",
                       "100GbE / InfiniBand networking",
                       "Qualcomm edge AI devices · 5G gateways · AI PCs",
                       "e& network: SIM, Toll Free 800, WhatsApp BSP, SMS",
                       "Physical e& devices: routers · IoT · cameras (future)"],
        "why": ("e&'s deepest moat. No AI competitor has sovereign UAE data centres, a "
                "telco network, SIM identity, and a physical device footprint. Everything "
                "above is software. This layer is physical."),
    },
]


FLOW_STEPS = [
    ("6", "Customers & Distribution", "SMB owner subscribes",
     "Selects a Customer Agent on the e& marketplace at AED 285 / month. One-click checkout on the e& account."),
    ("5", "Agentic Applications", "Customer Agent goes live",
     "Starts answering WhatsApp messages and phone calls in Arabic and English, books appointments, escalates to humans."),
    ("4B", "Forge · Development", "SDD blueprint reused",
     "Instantiated from a tested Customer Agent blueprint — prompts, guardrails, connectors and tests already proven."),
    ("4A", "Forge · Operations", "Control plane orchestrates",
     "Provisions the tenant, wires WhatsApp BSP + e& telephony, applies T1 compliance, starts metering token + GPU usage."),
    ("3", "Model Layer", "LLM — fixed or brokered",
     "Forge selects between Claude, GPT, Llama, or a sovereign model based on task type, latency, and data-residency rules."),
    ("2", "GPU Orchestration", "GPU orchestrator schedules",
     "Submits the inference job to the right compute target — H100/B200 on sovereign infra, with the right tenant quota."),
    ("1", "Infrastructure", "e& Data Centre runs the workload",
     "GPU runs in UAE sovereign DC. Network delivers the response in 1.2 seconds. Data never leaves UAE."),
    ("6", "Customers & Distribution", "e& invoice bills usage",
     "AED 285 added to the SMB's monthly e& bill. Revenue share computed automatically. One invoice for telco + AI."),
]


# ─────────────────────────────────────────────────────────────────────────────
# Helpers

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.shadow.inherit = False
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_w is not None:
            shape.line.width = line_w
    return shape


def add_text(slide, x, y, w, h, text, *,
             size=11, bold=False, color=BRAND["black"], font=FONT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             tracking=None, upper=False, line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text.upper() if upper else text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if tracking is not None:
        # python-pptx exposes character spacing via the spc attribute on rPr.
        from pptx.oxml.ns import qn
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(tracking))
    return tb


def add_badge(slide, x, y, text, *, bg=BRAND["red"], fg=BRAND["white"], size=9):
    width = Inches(0.05 * len(text) + 0.4)
    height = Inches(0.24)
    add_rect(slide, x, y, width, height, fill=bg)
    add_text(slide, x, y, width, height, text,
             size=size, bold=True, color=fg, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, upper=True, tracking=120)


def add_brand_stripe(slide, x, y, w):
    # Three-segment continuum stripe (rose → red → dark red), like the page hero.
    segs = [BRAND["red"], BRAND["red"], BRAND["red_dark"]]
    seg_w = w // len(segs)
    for i, c in enumerate(segs):
        add_rect(slide, x + seg_w * i, y, seg_w, Inches(0.05), fill=c)


def add_footer(slide, page_num):
    slide_w = Inches(13.333)
    slide_h = Inches(7.5)
    # Bottom divider
    add_rect(slide, Inches(0.5), slide_h - Inches(0.45), slide_w - Inches(1), Inches(0.01),
             fill=BRAND["border"])
    add_text(slide, Inches(0.5), slide_h - Inches(0.38), Inches(6), Inches(0.25),
             "e& AI Networks & Solutions × AIdeology  ·  AI Stack — six layers from silicon to customer",
             size=8.5, color=BRAND["grey"], tracking=80)
    add_text(slide, slide_w - Inches(1.2), slide_h - Inches(0.38), Inches(0.7), Inches(0.25),
             f"{page_num} / 10",
             size=8.5, bold=True, color=BRAND["grey_dark"], align=PP_ALIGN.RIGHT, tracking=80)


def add_section_title(slide, kicker, title, sub=None):
    add_brand_stripe(slide, Inches(0.5), Inches(0.5), Inches(2.2))
    add_text(slide, Inches(0.5), Inches(0.65), Inches(8), Inches(0.3),
             kicker, size=10, bold=True, color=BRAND["red"], upper=True, tracking=140)
    add_text(slide, Inches(0.5), Inches(0.95), Inches(12), Inches(0.7),
             title, size=28, bold=True, color=BRAND["black"], line_spacing=1.1)
    if sub:
        add_text(slide, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.6),
                 sub, size=12, color=BRAND["grey_mid"], line_spacing=1.45)


# ─────────────────────────────────────────────────────────────────────────────
# Slide builders

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])

    add_badge(s, Inches(0.5), Inches(0.5), "Article 4 · Technical Architecture")

    add_text(s, Inches(0.5), Inches(0.95), Inches(8), Inches(0.35),
             "e& AI Networks & Solutions × AIdeology · v0.1 · Confidential",
             size=10, color=BRAND["grey"], tracking=60)

    add_text(s, Inches(0.5), Inches(1.4), Inches(12), Inches(1.8),
             "The e& AI Stack",
             size=54, bold=True, color=BRAND["black"], line_spacing=1.05)

    add_text(s, Inches(0.5), Inches(2.55), Inches(12), Inches(0.8),
             "Six layers from silicon to customer.",
             size=24, color=BRAND["red"], bold=True, line_spacing=1.2)

    add_rect(s, Inches(0.5), Inches(3.35), Inches(2.5), Inches(0.04), fill=BRAND["red"])

    add_text(s, Inches(0.5), Inches(3.6), Inches(10.5), Inches(1.6),
             ("One AI stack across six layers. Forge sits in the middle as the AI control "
              "plane that operates the platform and the development layer that builds new "
              "AI fast. Below Forge: the GPU orchestrator on top of e&'s sovereign "
              "infrastructure. Above Forge: the agents customers actually buy, distributed "
              "through e&'s marketplace and invoice."),
             size=14, color=BRAND["grey_dark"], line_spacing=1.55)

    # Stat strip
    stats = [
        ("6 layers", "From silicon to customer"),
        ("Forge", "Sits at the centre"),
        ("4 of 6 layers", "Owned by e&"),
        ("240M+", "Direct billing relationships"),
    ]
    x0 = Inches(0.5)
    y0 = Inches(5.7)
    col_w = Inches(3.0)
    for i, (v, l) in enumerate(stats):
        add_text(s, x0 + col_w * i, y0, col_w, Inches(0.55),
                 v, size=24, bold=True, color=BRAND["red"])
        add_text(s, x0 + col_w * i, y0 + Inches(0.55), col_w, Inches(0.35),
                 l, size=9, bold=True, color=BRAND["grey"], upper=True, tracking=120)

    add_footer(s, 1)


def _draw_layer_strip(slide, x, y, w, layer, h=Inches(0.7)):
    """Compact horizontal layer card used in overview slides."""
    box_w = Inches(0.95)
    add_rect(slide, x, y, box_w, h, fill=layer["color"])
    add_text(slide, x, y, box_w, Inches(0.2), "LAYER",
             size=7.5, bold=True, color=BRAND["white"], align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.TOP, tracking=140)
    add_text(slide, x, y + Inches(0.2), box_w, h - Inches(0.2),
             layer["n"], size=24, bold=True, color=BRAND["white"],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    body_x = x + box_w
    body_w = w - box_w
    add_rect(slide, body_x, y, body_w, h, fill=BRAND["white"], line=BRAND["border"])
    add_text(slide, body_x + Inches(0.18), y + Inches(0.08), body_w - Inches(0.3), Inches(0.18),
             layer["tag"], size=8, bold=True, color=layer["color"], upper=True, tracking=120)
    add_text(slide, body_x + Inches(0.18), y + Inches(0.25), body_w - Inches(2.3), Inches(0.3),
             layer["name"], size=13, bold=True, color=BRAND["black"])
    add_text(slide, body_x + Inches(0.18), y + Inches(0.5), body_w - Inches(2.3), Inches(0.2),
             layer.get("owner", ""), size=9, color=BRAND["grey_mid"])
    # Right-side hint
    label = "FORGE · PLATFORM IP" if layer.get("is_forge") else "STACK LAYER"
    add_text(slide, body_x + body_w - Inches(2.1), y + Inches(0.25), Inches(2), Inches(0.25),
             label, size=8.5, bold=True, color=BRAND["grey"], upper=True, tracking=120,
             align=PP_ALIGN.RIGHT)


def slide_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        s, "Article 4 · Section 4.1",
        "The six-layer stack at a glance",
        ("Six layers, one platform. The arrows of value run top-down: customers consume "
         "agents; agents run on Forge; Forge brokers models; models execute on GPUs; GPUs "
         "live in e&'s sovereign infrastructure."))

    y = Inches(2.45)
    for L in STACK_LAYERS:
        _draw_layer_strip(s, Inches(0.5), y, Inches(12.3), L, h=Inches(0.68))
        y += Inches(0.74)

    add_footer(s, 2)


def slide_standard_layer(prs, layer, page_num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        s, f"Layer {layer['n']} · {layer['tag']}",
        layer["name"])

    # Left side: colored layer block
    add_rect(s, Inches(0.5), Inches(2.5), Inches(2.5), Inches(4.3), fill=layer["color"])
    add_text(s, Inches(0.5), Inches(2.6), Inches(2.5), Inches(0.3),
             "LAYER", size=10, bold=True, color=BRAND["white"],
             align=PP_ALIGN.CENTER, upper=True, tracking=160)
    add_text(s, Inches(0.5), Inches(2.9), Inches(2.5), Inches(1.6),
             layer["n"], size=110, bold=True, color=BRAND["white"],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    # Owner panel on left bottom
    add_text(s, Inches(0.5), Inches(5.6), Inches(2.5), Inches(0.25),
             "OWNER", size=9, bold=True, color=BRAND["white"],
             align=PP_ALIGN.CENTER, upper=True, tracking=140)
    add_text(s, Inches(0.7), Inches(5.85), Inches(2.1), Inches(0.85),
             layer["owner"], size=11, bold=True, color=BRAND["white"],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.35)

    # Right side: summary + components + why
    rx = Inches(3.2)
    rw = Inches(9.6)

    add_text(s, rx, Inches(2.5), rw, Inches(1.4),
             layer["summary"], size=12.5, color=BRAND["grey_dark"], line_spacing=1.55)

    # Components header
    add_text(s, rx, Inches(4.05), rw, Inches(0.25),
             "KEY COMPONENTS", size=9, bold=True, color=BRAND["grey"],
             upper=True, tracking=140)

    # Components 2-column grid
    comps = layer["components"]
    half = (len(comps) + 1) // 2
    col_w = rw / 2 - Inches(0.1)
    for col in range(2):
        cx = rx + col * (col_w + Inches(0.2))
        items = comps[col * half:(col + 1) * half]
        for i, c in enumerate(items):
            row_y = Inches(4.32) + Inches(0.26) * i
            add_rect(s, cx, row_y + Inches(0.09), Inches(0.06), Inches(0.06), fill=layer["color"])
            add_text(s, cx + Inches(0.16), row_y, col_w - Inches(0.16), Inches(0.26),
                     c, size=10.5, color=BRAND["grey_dark"], line_spacing=1.3)

    # Why it matters strip
    why_y = Inches(6.25)
    add_rect(s, rx, why_y, rw, Inches(0.04), fill=BRAND["border"])
    add_text(s, rx, why_y + Inches(0.12), rw, Inches(0.22),
             "WHY IT MATTERS", size=9, bold=True, color=layer["color"],
             upper=True, tracking=140)
    add_text(s, rx, why_y + Inches(0.38), rw, Inches(0.7),
             layer["why"], size=10.5, color=BRAND["grey_dark"], line_spacing=1.5)

    add_footer(s, page_num)


def slide_forge(prs, layer, page_num):
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # Section title (rendered manually to keep room for the dark Forge hero)
    add_brand_stripe(s, Inches(0.5), Inches(0.5), Inches(2.2))
    add_text(s, Inches(0.5), Inches(0.65), Inches(8), Inches(0.3),
             f"Layer {layer['n']} · {layer['tag']}",
             size=10, bold=True, color=BRAND["red"], upper=True, tracking=140)
    add_text(s, Inches(0.5), Inches(0.95), Inches(12), Inches(0.7),
             layer["name"],
             size=24, bold=True, color=BRAND["black"], line_spacing=1.1)

    # Dark Forge hero block
    hero_x = Inches(0.5)
    hero_y = Inches(1.85)
    hero_w = Inches(12.33)
    hero_h = Inches(1.55)
    add_rect(s, hero_x, hero_y, hero_w, hero_h, fill=BRAND["forge_a"])
    add_rect(s, hero_x, hero_y, hero_w, Inches(0.05), fill=BRAND["red"])  # top stripe

    # Large layer number
    add_rect(s, hero_x + Inches(0.4), hero_y + Inches(0.2), Inches(1.4), Inches(1.15),
             fill=BRAND["forge_b"])
    add_text(s, hero_x + Inches(0.4), hero_y + Inches(0.25), Inches(1.4), Inches(0.2),
             "LAYER", size=8.5, bold=True, color=BRAND["white"],
             align=PP_ALIGN.CENTER, upper=True, tracking=140)
    add_text(s, hero_x + Inches(0.4), hero_y + Inches(0.4), Inches(1.4), Inches(0.95),
             layer["n"], size=42, bold=True, color=BRAND["white"],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    # Forge tag + name + summary
    text_x = hero_x + Inches(2.1)
    text_w = hero_w - Inches(2.2)
    add_text(s, text_x, hero_y + Inches(0.18), text_w, Inches(0.25),
             layer["tag"], size=9.5, bold=True, color=BRAND["red"],
             upper=True, tracking=140)
    add_text(s, text_x, hero_y + Inches(0.42), text_w, Inches(0.45),
             layer["name"], size=20, bold=True, color=BRAND["white"], line_spacing=1.15)
    add_text(s, text_x, hero_y + Inches(0.92), text_w, Inches(0.6),
             layer["summary"], size=11, color=RGBColor(0xCC, 0xCC, 0xCC),
             line_spacing=1.5)

    # Two pillars
    pillar_y = hero_y + hero_h + Inches(0.2)
    pillar_h = Inches(4.05)
    pillar_w = (hero_w - Inches(0.2)) / 2
    for i, P in enumerate(layer["pillars"]):
        px = hero_x + (pillar_w + Inches(0.2)) * i
        bg = BRAND["panel_bg"] if i == 0 else BRAND["white"]
        add_rect(s, px, pillar_y, pillar_w, pillar_h, fill=bg, line=BRAND["border"])
        # Top accent
        add_rect(s, px, pillar_y, pillar_w, Inches(0.05), fill=P["color"])

        # Header strip
        add_rect(s, px, pillar_y + Inches(0.05), Inches(0.35), Inches(0.35), fill=P["color"])
        add_text(s, px, pillar_y + Inches(0.05), Inches(0.35), Inches(0.35),
                 "A" if i == 0 else "B", size=14, bold=True, color=BRAND["white"],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, px + Inches(0.45), pillar_y + Inches(0.06), pillar_w - Inches(0.5), Inches(0.2),
                 P["sub"], size=8.5, bold=True, color=P["color"], upper=True, tracking=140)
        add_text(s, px + Inches(0.45), pillar_y + Inches(0.22), pillar_w - Inches(0.5), Inches(0.28),
                 P["title"], size=13, bold=True, color=BRAND["black"])

        # Items
        for j, (t, d) in enumerate(P["items"]):
            row_y = pillar_y + Inches(0.6) + Inches(0.52) * j
            add_text(s, px + Inches(0.15), row_y, Inches(2.0), Inches(0.22),
                     t, size=10, bold=True, color=BRAND["black"])
            add_text(s, px + Inches(0.15), row_y + Inches(0.22), pillar_w - Inches(0.3), Inches(0.28),
                     d, size=8.5, color=BRAND["grey_mid"], line_spacing=1.35)
            if j < len(P["items"]) - 1:
                add_rect(s, px + Inches(0.15), row_y + Inches(0.5),
                         pillar_w - Inches(0.3), Inches(0.01), fill=BRAND["border"])

    add_footer(s, page_num)


def slide_flow(prs, page_num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        s, "Article 4 · Section 4.3",
        "End-to-end — one customer interaction across all six layers",
        "An SMB subscribes to a Customer Agent. A single interaction traverses every layer of the stack.")

    table_x = Inches(0.5)
    table_y = Inches(2.45)
    table_w = Inches(12.33)
    row_h = Inches(0.58)

    for i, (l, where, who, what) in enumerate(FLOW_STEPS):
        y = table_y + row_h * i
        # Layer pill
        is_forge = l.startswith("4")
        pill_bg = BRAND["red"] if is_forge else BRAND["black"]
        add_rect(s, table_x, y, Inches(0.85), row_h, fill=pill_bg)
        add_text(s, table_x, y + Inches(0.05), Inches(0.85), Inches(0.18),
                 "L", size=7.5, bold=True, color=BRAND["white"],
                 align=PP_ALIGN.CENTER, upper=True, tracking=140)
        add_text(s, table_x, y + Inches(0.18), Inches(0.85), row_h - Inches(0.18),
                 l, size=16 if len(l) > 1 else 18, bold=True, color=BRAND["white"],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Where + who cell
        wx = table_x + Inches(0.85)
        ww = Inches(3.2)
        add_rect(s, wx, y, ww, row_h, fill=BRAND["white"], line=BRAND["border"])
        add_text(s, wx + Inches(0.15), y + Inches(0.08), ww - Inches(0.3), Inches(0.2),
                 where, size=8.5, bold=True, color=BRAND["grey"], upper=True, tracking=120)
        add_text(s, wx + Inches(0.15), y + Inches(0.27), ww - Inches(0.3), Inches(0.3),
                 who, size=11, bold=True, color=BRAND["black"])

        # Description cell
        dx = wx + ww
        dw = table_w - Inches(0.85) - ww
        add_rect(s, dx, y, dw, row_h, fill=BRAND["white"], line=BRAND["border"])
        add_text(s, dx + Inches(0.2), y + Inches(0.1), dw - Inches(0.4), row_h - Inches(0.2),
                 what, size=10, color=BRAND["grey_dark"],
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)

    add_footer(s, page_num)


def slide_ownership(prs, page_num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        s, "Article 4 · Section 4.4",
        "Ownership map — what e& owns, what AIdeology builds",
        "Four of six layers transfer to e& ownership. Forge sits at the centre — AIdeology's licensed platform IP.")

    headers = ["Layer", "Name", "Owner", "What it means"]
    col_x = [Inches(0.5), Inches(1.5), Inches(5.4), Inches(8.4)]
    col_w = [Inches(1.0), Inches(3.9), Inches(3.0), Inches(4.4)]

    # Header row
    hy = Inches(2.4)
    add_rect(s, col_x[0], hy, sum(col_w, Inches(0)), Inches(0.42), fill=BRAND["lightgrey"])
    for i, h in enumerate(headers):
        add_text(s, col_x[i] + Inches(0.15), hy + Inches(0.1), col_w[i] - Inches(0.2), Inches(0.3),
                 h, size=9, bold=True, color=BRAND["grey"], upper=True, tracking=140)

    rows = [
        ("6", "Customers & Distribution", "e& (100%)", "Customer relationship, brand, billing, support"),
        ("5", "Agentic Applications",     "e& (transfer)", "Agent IP transfers wave-by-wave under build-then-transfer"),
        ("4", "Forge — Control Plane + Dev", "AIdeology (licensed to e&)", "Durable platform IP; perpetual non-exclusive licence to e&"),
        ("3", "Model Layer",              "Mixed", "Third-party APIs + e&-hosted sovereign Arabic models"),
        ("2", "GPU Orchestration",        "e& (Open Innovation)", "Sovereign GPU orchestrator on e& compute estate"),
        ("1", "Infrastructure",           "e& (100%)", "Data centres, network, hardware, sovereign devices"),
    ]
    for i, row in enumerate(rows):
        ry = hy + Inches(0.42) + Inches(0.55) * i
        is_forge = row[0] == "4"
        bg = BRAND["rose_bg"] if is_forge else BRAND["white"]
        add_rect(s, col_x[0], ry, sum(col_w, Inches(0)), Inches(0.55),
                 fill=bg, line=BRAND["border"])
        add_text(s, col_x[0] + Inches(0.15), ry + Inches(0.15), col_w[0] - Inches(0.2), Inches(0.3),
                 row[0], size=14, bold=True, color=BRAND["red"])
        add_text(s, col_x[1] + Inches(0.15), ry + Inches(0.17), col_w[1] - Inches(0.2), Inches(0.3),
                 row[1], size=11, bold=True, color=BRAND["black"])
        add_text(s, col_x[2] + Inches(0.15), ry + Inches(0.17), col_w[2] - Inches(0.2), Inches(0.3),
                 row[2], size=10.5, color=BRAND["grey_dark"])
        add_text(s, col_x[3] + Inches(0.15), ry + Inches(0.17), col_w[3] - Inches(0.2), Inches(0.3),
                 row[3], size=10, color=BRAND["grey_mid"])

    # Bottom takeaway band
    by = Inches(6.55)
    add_rect(s, Inches(0.5), by, Inches(12.33), Inches(0.7), fill=BRAND["black"])
    add_rect(s, Inches(0.5), by, Inches(0.1), Inches(0.7), fill=BRAND["red"])
    add_text(s, Inches(0.8), by + Inches(0.1), Inches(11.5), Inches(0.25),
             "THE BOTTOM LINE", size=9, bold=True, color=BRAND["red"],
             upper=True, tracking=140)
    add_text(s, Inches(0.8), by + Inches(0.32), Inches(11.5), Inches(0.32),
             "e& owns the customer, the agents, the GPUs and the data centres. AIdeology builds and licenses Forge — the durable platform IP that runs the whole stack.",
             size=11, color=BRAND["white"], line_spacing=1.4)

    add_footer(s, page_num)


# ─────────────────────────────────────────────────────────────────────────────
# Build presentation

def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 — Cover
    slide_cover(prs)
    # 2 — Overview of six layers
    slide_overview(prs)
    # 3–8 — One slide per layer (6 → 1)
    pn = 3
    for L in STACK_LAYERS:
        if L.get("is_forge"):
            slide_forge(prs, L, pn)
        else:
            slide_standard_layer(prs, L, pn)
        pn += 1
    # 9 — End-to-end flow
    slide_flow(prs, pn)
    pn += 1
    # 10 — Ownership map / close
    slide_ownership(prs, pn)

    out = "e&_AI_Stack_Six_Layers.pptx"
    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    build()
