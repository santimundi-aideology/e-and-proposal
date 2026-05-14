"""Professional Services for GPUaaS — commercial-operating boundary deck.

Audience: Fikreab (Core42 partnership commercial deck), e& stakeholders.
Goal: unblock the GPUaaS commercial narrative with clear ownership, services
and pricing posture between AIdeology, e&, and Core42.

Reuses brand helpers from generate_stack_deck.py.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from generate_stack_deck import (
    BRAND, FONT,
    add_rect, add_text, add_badge, add_brand_stripe,
)


# ─────────────────────────────────────────────────────────────────────────────
# Local footer (override to retitle for this deck)

def add_footer(slide, page_num, total=9):
    slide_w = Inches(13.333)
    slide_h = Inches(7.5)
    add_rect(slide, Inches(0.5), slide_h - Inches(0.45), slide_w - Inches(1), Inches(0.01),
             fill=BRAND["border"])
    add_text(slide, Inches(0.5), slide_h - Inches(0.38), Inches(8), Inches(0.25),
             "AIdeology × e& × Core42  ·  Professional Services for GPUaaS  ·  Boundary & commercial pack",
             size=8.5, color=BRAND["grey"], tracking=80)
    add_text(slide, slide_w - Inches(1.2), slide_h - Inches(0.38), Inches(0.7), Inches(0.25),
             f"{page_num} / {total}",
             size=8.5, bold=True, color=BRAND["grey_dark"], align=PP_ALIGN.RIGHT, tracking=80)


def add_section_title(slide, kicker, title, sub=None):
    add_brand_stripe(slide, Inches(0.5), Inches(0.5), Inches(2.2))
    add_text(slide, Inches(0.5), Inches(0.65), Inches(10), Inches(0.3),
             kicker, size=10, bold=True, color=BRAND["red"], upper=True, tracking=140)
    add_text(slide, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.7),
             title, size=26, bold=True, color=BRAND["black"], line_spacing=1.1)
    if sub:
        add_text(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.6),
                 sub, size=12, color=BRAND["grey_mid"], line_spacing=1.45)


# Three-party color coding used throughout
PARTY_COLORS = {
    "A": BRAND["red"],          # AIdeology
    "E": RGBColor(0x18, 0x5F, 0xA5),  # e&
    "C": RGBColor(0x00, 0x4B, 0x2E),  # Core42
}
PARTY_NAMES = {"A": "AIdeology", "E": "e&", "C": "Core42"}


def party_chip(slide, x, y, w, h, code):
    add_rect(slide, x, y, w, h, fill=PARTY_COLORS[code])
    add_text(slide, x, y, w, h, code,
             size=11, bold=True, color=BRAND["white"],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ─────────────────────────────────────────────────────────────────────────────
# Slide builders

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])

    add_badge(s, Inches(0.5), Inches(0.5), "Internal · Confidential")

    add_text(s, Inches(0.5), Inches(0.95), Inches(8), Inches(0.35),
             "Prepared for Fikreab's GPUaaS commercial deck · May 2026",
             size=10, color=BRAND["grey"], tracking=60)

    add_text(s, Inches(0.5), Inches(1.4), Inches(12), Inches(1.8),
             "Professional Services for GPUaaS",
             size=48, bold=True, color=BRAND["black"], line_spacing=1.05)

    add_text(s, Inches(0.5), Inches(2.45), Inches(12), Inches(0.7),
             "Scope, boundaries, and commercial model.",
             size=22, color=BRAND["red"], bold=True, line_spacing=1.2)

    add_rect(s, Inches(0.5), Inches(3.2), Inches(2.5), Inches(0.04), fill=BRAND["red"])

    add_text(s, Inches(0.5), Inches(3.45), Inches(11), Inches(1.6),
             ("One pack to unblock the commercial narrative for GPUaaS — what "
              "AIdeology, e&, and Core42 each own across professional services, "
              "where the overlaps are, and how each party earns. Built to be "
              "decided, not debated."),
             size=14, color=BRAND["grey_dark"], line_spacing=1.55)

    # Three-party block
    parties = [
        ("A", "AIdeology", "AI plane, agentic delivery, MLOps, HPC RA design"),
        ("E", "e&", "Commercial front, customer ownership, managed service wrap"),
        ("C", "Core42", "Sovereign cloud / GPU infra hosting, capacity layer"),
    ]
    x0 = Inches(0.5); y0 = Inches(5.4)
    col_w = Inches(4.1)
    for i, (code, name, desc) in enumerate(parties):
        cx = x0 + col_w * i + Inches(0.05 if i > 0 else 0)
        add_rect(s, cx, y0, col_w - Inches(0.1), Inches(1.5),
                 fill=BRAND["white"], line=BRAND["border"])
        add_rect(s, cx, y0, Inches(0.18), Inches(1.5), fill=PARTY_COLORS[code])
        add_text(s, cx + Inches(0.32), y0 + Inches(0.15), col_w - Inches(0.5), Inches(0.3),
                 PARTY_NAMES[code], size=9, bold=True, color=PARTY_COLORS[code],
                 upper=True, tracking=140)
        add_text(s, cx + Inches(0.32), y0 + Inches(0.4), col_w - Inches(0.5), Inches(0.4),
                 name, size=18, bold=True, color=BRAND["black"])
        add_text(s, cx + Inches(0.32), y0 + Inches(0.85), col_w - Inches(0.5), Inches(0.6),
                 desc, size=11, color=BRAND["grey_mid"], line_spacing=1.4)

    add_footer(s, 1)


def slide_objective(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        s, "Slide 02 · Objective",
        "Why this pack exists — and the decisions it asks for",
        ("Fikreab needs a clean GPUaaS commercial story. This pack provides the "
         "professional services scope, ownership boundaries, and commercial "
         "frame so the Core42 narrative can land without ambiguity."))

    # Two columns: "What this is" / "What this is not"
    col_w = Inches(5.9)
    y0 = Inches(2.6)
    for i, (title, color, items) in enumerate([
        ("WHAT THIS PACK IS", BRAND["red"], [
            "Commercial-operating boundary deck for GPUaaS",
            "Clear ownership map: AIdeology / e& / Core42",
            "Professional services catalogue with pricing posture",
            "Overlap zones with single-accountable-owner rules",
            "Input deck for Fikreab's Core42 commercial narrative",
        ]),
        ("WHAT IT IS NOT", BRAND["grey"], [
            "Not a technical architecture deep-dive",
            "Not a final SoW or binding pricing",
            "Not a Core42 product / capacity proposal",
            "Not a customer-facing slide deck — internal alignment",
            "Not a substitute for a co-sell agreement",
        ]),
    ]):
        x = Inches(0.5) + (col_w + Inches(0.3)) * i
        add_rect(s, x, y0, col_w, Inches(2.55), fill=BRAND["white"], line=BRAND["border"])
        add_rect(s, x, y0, col_w, Inches(0.06), fill=color)
        add_text(s, x + Inches(0.3), y0 + Inches(0.2), col_w - Inches(0.5), Inches(0.3),
                 title, size=10, bold=True, color=color, upper=True, tracking=140)
        for j, it in enumerate(items):
            iy = y0 + Inches(0.6) + Inches(0.35) * j
            add_rect(s, x + Inches(0.3), iy + Inches(0.1), Inches(0.06), Inches(0.06), fill=color)
            add_text(s, x + Inches(0.45), iy, col_w - Inches(0.55), Inches(0.32),
                     it, size=11.5, color=BRAND["grey_dark"], line_spacing=1.35)

    # Bottom: 3 decisions needed
    by = Inches(5.4)
    add_text(s, Inches(0.5), by, Inches(12), Inches(0.3),
             "3 DECISIONS NEEDED TODAY", size=10, bold=True, color=BRAND["red"],
             upper=True, tracking=160)

    decisions = [
        ("01", "Ownership boundaries", "Approve the AIdeology / e& / Core42 RACI on professional services."),
        ("02", "Commercial frame",     "Lock the three commercial packages: fixed build, advisory retainer, managed."),
        ("03", "Overlap rules",        "Approve the single-owner rule for the 4 overlap zones."),
    ]
    dy = Inches(5.75)
    dw = Inches(4.0)
    for i, (n, t, d) in enumerate(decisions):
        dx = Inches(0.5) + (dw + Inches(0.15)) * i
        add_rect(s, dx, dy, dw, Inches(1.3), fill=BRAND["lightgrey"], line=BRAND["border"])
        add_rect(s, dx, dy, Inches(0.5), Inches(1.3), fill=BRAND["red"])
        add_text(s, dx, dy + Inches(0.32), Inches(0.5), Inches(0.6),
                 n, size=22, bold=True, color=BRAND["white"],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, dx + Inches(0.62), dy + Inches(0.18), dw - Inches(0.7), Inches(0.3),
                 t, size=12.5, bold=True, color=BRAND["black"])
        add_text(s, dx + Inches(0.62), dy + Inches(0.48), dw - Inches(0.7), Inches(0.7),
                 d, size=10.5, color=BRAND["grey_mid"], line_spacing=1.45)

    add_footer(s, 2)


def slide_scope_map(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        s, "Slide 03 · End-to-end scope",
        "From silicon to live agentic solution",
        ("Eight layers of professional services across four phases. AIdeology "
         "leads the AI / platform plane; Core42 provides sovereign GPU "
         "capacity; e& fronts the customer."))

    # Phase band header (4 phases)
    phases = [
        ("DESIGN",  "Wks 0–8",  BRAND["red"]),
        ("BUILD",   "Wks 6–20", RGBColor(0x7A, 0x52, 0xF4)),
        ("LAUNCH",  "Wks 18–26", RGBColor(0x18, 0x5F, 0xA5)),
        ("OPERATE", "Wk 24+",   RGBColor(0x0F, 0x6B, 0x49)),
    ]
    pw = Inches(3.0)
    py = Inches(2.5)
    for i, (n, t, c) in enumerate(phases):
        px = Inches(0.5) + (pw + Inches(0.08)) * i
        add_rect(s, px, py, pw, Inches(0.55), fill=c)
        add_text(s, px, py, pw, Inches(0.28),
                 n, size=11, bold=True, color=BRAND["white"],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, tracking=180,
                 line_spacing=1.1)
        add_text(s, px, py + Inches(0.28), pw, Inches(0.27),
                 t, size=9, color=BRAND["white"],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, tracking=140)

    # 8 layers (bottom-up build, but rendered top-down here for reading flow)
    layers = [
        ("08", "Customer Outcome",          "Live agentic solution on the customer's workflow · billed on e&", "A · E"),
        ("07", "Agentic Applications",      "Multi-agent solutions · contact centre · doc intelligence · approvals", "A"),
        ("06", "Forge AI Control Plane",    "Orchestration · LLM gateway · multi-tenancy · governance · billing", "A"),
        ("05", "MLOps & Model Layer",       "NVIDIA AI Enterprise · Fleet Command · serving · fine-tuning · drift", "A"),
        ("04", "Data Platform",             "Lustre · GPFS · GPUDirect · RAPIDS · governance · lineage", "A"),
        ("03", "Compute Orchestration",     "Kubernetes · Slurm · MIG · multi-tenant isolation · usage meters", "A · C"),
        ("02", "GPU Hardware & Fabric",     "NVIDIA DGX · HGX · NVL72 · RA design · BoM · procurement advisory", "A · C"),
        ("01", "Site & Facility Foundation","Power · cooling · fabric · physical security · sovereign DC", "E · C"),
    ]
    y0 = Inches(3.3)
    row_h = Inches(0.42)
    for i, (n, t, d, owner) in enumerate(layers):
        y = y0 + row_h * i
        # Number tab
        add_rect(s, Inches(0.5), y, Inches(0.55), row_h, fill=BRAND["black"])
        add_text(s, Inches(0.5), y, Inches(0.55), row_h,
                 n, size=14, bold=True, color=BRAND["white"],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Body
        add_rect(s, Inches(1.05), y, Inches(9.55), row_h,
                 fill=BRAND["white"], line=BRAND["border"])
        add_text(s, Inches(1.2), y + Inches(0.05), Inches(3.0), Inches(0.32),
                 t, size=11.5, bold=True, color=BRAND["black"],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.2), y + Inches(0.05), Inches(6.3), Inches(0.32),
                 d, size=10, color=BRAND["grey_mid"],
                 anchor=MSO_ANCHOR.MIDDLE)
        # Owner tag
        add_rect(s, Inches(10.6), y, Inches(2.23), row_h,
                 fill=BRAND["lightgrey"], line=BRAND["border"])
        add_text(s, Inches(10.7), y + Inches(0.05), Inches(2.0), Inches(0.32),
                 "OWNER", size=8, bold=True, color=BRAND["grey"],
                 upper=True, tracking=140, anchor=MSO_ANCHOR.TOP)
        add_text(s, Inches(10.7), y + Inches(0.16), Inches(2.0), Inches(0.25),
                 owner, size=10.5, bold=True, color=BRAND["red"],
                 anchor=MSO_ANCHOR.TOP)

    # Legend
    ly = Inches(6.85)
    add_text(s, Inches(0.5), ly, Inches(2), Inches(0.25),
             "LEGEND", size=9, bold=True, color=BRAND["grey"], tracking=160)
    legend = [("A", "AIdeology"), ("E", "e&"), ("C", "Core42")]
    lx = Inches(1.4)
    for code, name in legend:
        party_chip(s, lx, ly - Inches(0.04), Inches(0.3), Inches(0.3), code)
        add_text(s, lx + Inches(0.4), ly + Inches(0.02), Inches(1.4), Inches(0.25),
                 name, size=10.5, color=BRAND["grey_dark"], bold=True)
        lx += Inches(2.0)

    add_footer(s, 3)


def slide_raci(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        s, "Slide 04 · Boundary map (RACI)",
        "Who does what — AIdeology · e& · Core42",
        ("Workstream-level ownership. R = Responsible (does the work). "
         "A = Accountable. S = Supports. — = Not involved."))

    rows = [
        ("Strategy & architecture (AI plane)",  "R/A", "S", "S"),
        ("Reference architecture (HPC / GPU)",  "R/A", "S", "S"),
        ("GPU infra hosting & operations",      "S",  "S", "R/A"),
        ("Network, security & sovereignty",     "S",  "R/A", "S"),
        ("Platform build (Forge, MLOps, data)", "R/A", "S", "—"),
        ("Agentic solution build (SDD)",        "R/A", "S", "—"),
        ("Customer-facing managed service",     "S",  "R/A", "S"),
        ("L3/L4 AI platform support",            "R/A", "S", "—"),
        ("Commercial / GTM ownership",           "S",  "R/A", "S"),
        ("Compliance, audit & evidence packs",   "R",  "A", "S"),
    ]

    # Header
    hy = Inches(2.4)
    headers = ["Workstream", "AIdeology", "e&", "Core42"]
    col_x = [Inches(0.5), Inches(7.0), Inches(8.95), Inches(10.9)]
    col_w = [Inches(6.4), Inches(1.85), Inches(1.85), Inches(1.85)]

    add_rect(s, col_x[0], hy, sum(col_w, Inches(0)), Inches(0.42), fill=BRAND["lightgrey"])
    add_text(s, col_x[0] + Inches(0.2), hy + Inches(0.1), col_w[0] - Inches(0.2), Inches(0.3),
             headers[0], size=10, bold=True, color=BRAND["grey"], upper=True, tracking=140)
    for i in range(1, 4):
        add_text(s, col_x[i], hy + Inches(0.1), col_w[i], Inches(0.3),
                 headers[i], size=10, bold=True, color=BRAND["grey"],
                 upper=True, tracking=140, align=PP_ALIGN.CENTER)

    row_h = Inches(0.39)
    for i, row in enumerate(rows):
        ry = hy + Inches(0.42) + row_h * i
        bg = BRAND["white"] if i % 2 == 0 else BRAND["lightgrey"]
        add_rect(s, col_x[0], ry, sum(col_w, Inches(0)), row_h, fill=bg, line=BRAND["border"])
        add_text(s, col_x[0] + Inches(0.2), ry + Inches(0.08), col_w[0] - Inches(0.3), Inches(0.3),
                 row[0], size=11, bold=True, color=BRAND["black"],
                 anchor=MSO_ANCHOR.MIDDLE)
        for j, val in enumerate(row[1:], start=1):
            party = ["A", "E", "C"][j - 1]
            color = PARTY_COLORS[party] if val != "—" else BRAND["grey_soft"]
            # Render as pill
            cell_cx = col_x[j] + col_w[j] / 2
            pill_w = Inches(0.85)
            pill_h = Inches(0.26)
            px = cell_cx - pill_w / 2
            py_ = ry + (row_h - pill_h) / 2
            add_rect(s, px, py_, pill_w, pill_h,
                     fill=color if val != "—" else BRAND["white"],
                     line=color if val != "—" else BRAND["border"])
            add_text(s, px, py_, pill_w, pill_h, val,
                     size=10, bold=True,
                     color=BRAND["white"] if val != "—" else BRAND["grey_soft"],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom rule strip
    by = Inches(6.75)
    add_rect(s, Inches(0.5), by, Inches(12.33), Inches(0.42), fill=BRAND["black"])
    add_rect(s, Inches(0.5), by, Inches(0.08), Inches(0.42), fill=BRAND["red"])
    add_text(s, Inches(0.7), by + Inches(0.05), Inches(12), Inches(0.16),
             "GOLDEN RULE", size=8.5, bold=True, color=BRAND["red"],
             upper=True, tracking=160)
    add_text(s, Inches(0.7), by + Inches(0.2), Inches(12), Inches(0.25),
             "Every overlap resolves to a single Accountable owner. No shared accountability.",
             size=11, color=BRAND["white"])

    add_footer(s, 4)


def slide_overlap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        s, "Slide 05 · Overlap zones",
        "Where roles touch — and how we resolve them",
        ("Four overlap zones surface in every GPUaaS deal. Each has a single-"
         "accountable-owner rule so the customer never sees a finger-point."))

    zones = [
        ("01", "Infra architecture vs infra hosting",
         "AIdeology designs the GPU / fabric / storage RA. Core42 operates the hosted infra.",
         "AIdeology = Accountable for architecture decisions and BoM. Core42 = Accountable for hosting performance & uptime."),
        ("02", "Platform ops vs cloud ops",
         "Forge / MLOps / data plane vs cloud-resource ops.",
         "AIdeology = Accountable for platform ops (L3/L4 AI). Core42 = Accountable for cloud-side resource health."),
        ("03", "Security controls vs security operations",
         "Policy, RBAC, governance, audit fabric vs SOC and run-time security.",
         "AIdeology = Accountable for AI / data controls. e& (via Help AG) = Accountable for security operations."),
        ("04", "Customer-facing managed service vs back-end managed service",
         "Customer-visible SLA vs supplier-side platform SLA.",
         "e& = Accountable for customer SLA. AIdeology = Accountable for platform SLA underneath."),
    ]

    grid_x = Inches(0.5)
    grid_y = Inches(2.5)
    col_w = Inches(6.16)
    row_h = Inches(2.15)

    for i, (n, t, d, rule) in enumerate(zones):
        col = i % 2
        row = i // 2
        x = grid_x + (col_w + Inches(0.15)) * col
        y = grid_y + (row_h + Inches(0.15)) * row
        add_rect(s, x, y, col_w, row_h, fill=BRAND["white"], line=BRAND["border"])
        add_rect(s, x, y, col_w, Inches(0.06), fill=BRAND["red"])
        # Number + title
        add_text(s, x + Inches(0.3), y + Inches(0.22), Inches(0.5), Inches(0.45),
                 n, size=28, bold=True, color=BRAND["red"], font=FONT,
                 anchor=MSO_ANCHOR.TOP)
        add_text(s, x + Inches(1.0), y + Inches(0.3), col_w - Inches(1.2), Inches(0.45),
                 t, size=14, bold=True, color=BRAND["black"], line_spacing=1.2)
        # Body
        add_text(s, x + Inches(0.3), y + Inches(0.85), col_w - Inches(0.5), Inches(0.55),
                 d, size=11, color=BRAND["grey_mid"], line_spacing=1.45)
        # Rule strip
        add_rect(s, x + Inches(0.3), y + Inches(1.4), col_w - Inches(0.5), Inches(0.65),
                 fill=BRAND["lightgrey"])
        add_text(s, x + Inches(0.45), y + Inches(1.45), Inches(2), Inches(0.2),
                 "RESOLUTION", size=8.5, bold=True, color=BRAND["red"],
                 upper=True, tracking=160)
        add_text(s, x + Inches(0.45), y + Inches(1.65), col_w - Inches(0.8), Inches(0.4),
                 rule, size=10, color=BRAND["black"], line_spacing=1.4, bold=True)

    add_footer(s, 5)


def slide_catalogue(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        s, "Slide 06 · Services catalogue",
        "What AIdeology brings into the partnership",
        ("Four service families. Each is a fixed-fee or retainer engagement and "
         "can be invoked independently or as part of a programme."))

    families = [
        ("HPC & Infrastructure", BRAND["black"], [
            "GPU Infrastructure Assessment",
            "Reference Architecture Design (RTX PRO · HGX · NVL72)",
            "BoM Optimisation & Procurement Advisory",
            "Networking & Fabric Advisory (InfiniBand / Spectrum-X)",
            "100 MW DC Fulfilment Programme",
            "Performance Benchmarking (MLPerf-aligned)",
        ]),
        ("Platform & MLOps", RGBColor(0x18, 0x5F, 0xA5), [
            "Forge Control Plane Deployment",
            "LLM Gateway & Model Router",
            "MLOps Infrastructure (NVIDIA AI Enterprise · Fleet Command)",
            "Model Serving (vLLM · TGI · Triton)",
            "Fine-Tuning Pipelines (DeepSpeed · LoRA · QLoRA)",
            "AI Data Platform (Lustre / GPFS · GPUDirect · RAPIDS)",
        ]),
        ("Agentic AI Delivery", BRAND["red"], [
            "AI Strategy & AIdeation",
            "Spec-Driven Design (SDD) methodology",
            "Custom Agent Development (multi-agent solutions)",
            "Adapted Agent Uplift (vertical hardening)",
            "Voice & Arabic Multilingual NLP",
            "Robotics, Simulation & 3D AI (Omniverse · Isaac Sim)",
        ]),
        ("Operate & Evolve", RGBColor(0x0F, 0x6B, 0x49), [
            "Managed AI Operations (24/7 L3/L4)",
            "Sovereign Managed Service (customer-site)",
            "Model Monitoring & Drift Detection",
            "GPU Fleet Management & Capacity",
            "Centre of Excellence build & Knowledge Transfer",
            "Continuous Agent Evolution",
        ]),
    ]

    gx = Inches(0.5); gy = Inches(2.45)
    col_w = Inches(3.05)
    row_h = Inches(4.45)
    for i, (name, color, items) in enumerate(families):
        x = gx + (col_w + Inches(0.06)) * i
        add_rect(s, x, gy, col_w, row_h, fill=BRAND["white"], line=BRAND["border"])
        add_rect(s, x, gy, col_w, Inches(0.55), fill=color)
        add_text(s, x + Inches(0.2), gy + Inches(0.16), col_w - Inches(0.3), Inches(0.3),
                 name, size=12, bold=True, color=BRAND["white"])
        for j, it in enumerate(items):
            iy = gy + Inches(0.75) + Inches(0.6) * j
            add_rect(s, x + Inches(0.2), iy + Inches(0.05), Inches(0.22), Inches(0.22),
                     fill=color)
            add_text(s, x + Inches(0.2), iy + Inches(0.05), Inches(0.22), Inches(0.22),
                     f"{j+1:02d}", size=8, bold=True, color=BRAND["white"],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT)
            add_text(s, x + Inches(0.5), iy, col_w - Inches(0.7), Inches(0.55),
                     it, size=10, color=BRAND["black"], line_spacing=1.4,
                     anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 6)


def slide_commercial(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        s, "Slide 07 · Commercial frame",
        "Three packaging blocks · how each party earns",
        ("Modular commercial structure. Customers can engage any block "
         "independently. Pricing is indicative — final terms per SoW."))

    blocks = [
        {
            "name": "DESIGN & BUILD",
            "color": BRAND["red"],
            "model": "Fixed-fee · milestone-based",
            "scope": ["RA design", "Platform build (Forge / MLOps / data)",
                      "Agent build (SDD)", "Site deployment support"],
            "price": "$300K – $1M",
            "unit":  "per programme",
            "earn":  ("AIdeology invoices fixed milestones. e& adds wrap fee. "
                      "Core42 invoices its own cloud / hosting set-up."),
        },
        {
            "name": "ADVISORY RETAINER",
            "color": RGBColor(0x18, 0x5F, 0xA5),
            "model": "Monthly retainer",
            "scope": ["2–3 senior architects on-account",
                      "Capacity & roadmap reviews",
                      "Vendor coordination",
                      "Procurement advisory"],
            "price": "$40K – $60K",
            "unit":  "per month",
            "earn":  ("AIdeology bills monthly. e& adds margin into customer "
                      "wrap. Core42 not directly involved unless infra advice."),
        },
        {
            "name": "MANAGED SERVICE",
            "color": RGBColor(0x0F, 0x6B, 0x49),
            "model": "Retainer + SLA",
            "scope": ["24/7 L3/L4 platform support",
                      "Model monitoring & drift",
                      "Continuous tuning",
                      "Sovereign managed option"],
            "price": "$25K – $60K",
            "unit":  "per month",
            "earn":  ("e& fronts customer SLA. AIdeology underwrites platform "
                      "SLA. Core42 contributes infra SLA where it hosts."),
        },
    ]

    gx = Inches(0.5); gy = Inches(2.45)
    col_w = Inches(4.16)
    row_h = Inches(4.5)
    for i, b in enumerate(blocks):
        x = gx + (col_w + Inches(0.06)) * i
        add_rect(s, x, gy, col_w, row_h, fill=BRAND["white"], line=BRAND["border"])
        # Header band
        add_rect(s, x, gy, col_w, Inches(0.75), fill=b["color"])
        add_text(s, x + Inches(0.25), gy + Inches(0.13), col_w - Inches(0.5), Inches(0.32),
                 b["name"], size=13, bold=True, color=BRAND["white"], tracking=120)
        add_text(s, x + Inches(0.25), gy + Inches(0.45), col_w - Inches(0.5), Inches(0.27),
                 b["model"], size=10, color=BRAND["white"], tracking=80)
        # Price
        py = gy + Inches(0.85)
        add_text(s, x + Inches(0.25), py, col_w - Inches(0.5), Inches(0.55),
                 b["price"], size=26, bold=True, color=b["color"], font=FONT)
        add_text(s, x + Inches(0.25), py + Inches(0.5), col_w - Inches(0.5), Inches(0.25),
                 b["unit"], size=10, color=BRAND["grey"], upper=True, tracking=140)
        # Divider
        add_rect(s, x + Inches(0.25), gy + Inches(1.7), col_w - Inches(0.5), Inches(0.01),
                 fill=BRAND["border"])
        # Scope items
        add_text(s, x + Inches(0.25), gy + Inches(1.8), col_w - Inches(0.5), Inches(0.22),
                 "SCOPE", size=8.5, bold=True, color=BRAND["grey"], tracking=160)
        for j, it in enumerate(b["scope"]):
            iy = gy + Inches(2.05) + Inches(0.3) * j
            add_rect(s, x + Inches(0.25), iy + Inches(0.1), Inches(0.07), Inches(0.07),
                     fill=b["color"])
            add_text(s, x + Inches(0.42), iy, col_w - Inches(0.6), Inches(0.3),
                     it, size=10.5, color=BRAND["grey_dark"], line_spacing=1.4)
        # Earn footer
        ey = gy + row_h - Inches(1.1)
        add_rect(s, x + Inches(0.15), ey, col_w - Inches(0.3), Inches(0.95),
                 fill=BRAND["lightgrey"])
        add_text(s, x + Inches(0.3), ey + Inches(0.1), col_w - Inches(0.6), Inches(0.2),
                 "WHO EARNS", size=8.5, bold=True, color=BRAND["red"],
                 upper=True, tracking=160)
        add_text(s, x + Inches(0.3), ey + Inches(0.32), col_w - Inches(0.6), Inches(0.6),
                 b["earn"], size=10, color=BRAND["grey_dark"], line_spacing=1.4)

    # Footer note
    add_text(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.2),
             "All blocks structured around milestone acceptance · 15% mobilisation advance · "
             "net 30 days from certification.",
             size=10, color=BRAND["grey"], align=PP_ALIGN.CENTER)

    add_footer(s, 7)


def slide_timeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        s, "Slide 08 · Delivery model",
        "30 · 60 · 90 day view, with 6-month landing",
        ("Stage-gated delivery. Each gate has a named go/no-go criterion. "
         "Parallel-loaded where dependencies allow."))

    stages = [
        ("DAY 0–30",  "DESIGN LOCK",
         ["AIdeation & use-case backlog",
          "Reference Architecture sign-off",
          "Commercial frame agreed",
          "Core42 cloud capacity confirmed"]),
        ("DAY 31–60", "BUILD KICK-OFF",
         ["Hardware procurement underway",
          "Platform build sprint 0",
          "Agent SDD specs signed",
          "Site readiness validated"]),
        ("DAY 61–90", "BETA LIVE",
         ["First agent in closed beta",
          "Platform deployed in dev/staging",
          "GPU pod commissioning started",
          "Security & compliance review passed"]),
        ("MONTH 4–6", "GA & OPERATE",
         ["Production launch on sovereign infra",
          "Managed service contract active",
          "Customer SLA dashboards live",
          "Continuous evolution backlog opened"]),
    ]

    gx = Inches(0.5); gy = Inches(2.6)
    col_w = Inches(3.13)
    row_h = Inches(3.6)
    colors = [BRAND["red"], RGBColor(0x7A, 0x52, 0xF4),
              RGBColor(0x18, 0x5F, 0xA5), RGBColor(0x0F, 0x6B, 0x49)]

    for i, (period, gate, items) in enumerate(stages):
        x = gx + (col_w + Inches(0.05)) * i
        color = colors[i]
        add_rect(s, x, gy, col_w, row_h, fill=BRAND["white"], line=BRAND["border"])
        add_rect(s, x, gy, col_w, Inches(0.06), fill=color)
        # Period band
        add_rect(s, x, gy + Inches(0.06), col_w, Inches(0.42), fill=BRAND["black"])
        add_text(s, x + Inches(0.2), gy + Inches(0.13), col_w - Inches(0.3), Inches(0.3),
                 period, size=12, bold=True, color=BRAND["white"], tracking=160)
        # Gate badge
        gy2 = gy + Inches(0.6)
        add_text(s, x + Inches(0.2), gy2, col_w - Inches(0.3), Inches(0.24),
                 "GATE", size=8.5, bold=True, color=color, tracking=180)
        add_text(s, x + Inches(0.2), gy2 + Inches(0.22), col_w - Inches(0.3), Inches(0.4),
                 gate, size=16, bold=True, color=BRAND["black"], line_spacing=1.15)
        # Items
        for j, it in enumerate(items):
            iy = gy + Inches(1.45) + Inches(0.45) * j
            add_rect(s, x + Inches(0.2), iy + Inches(0.08), Inches(0.06), Inches(0.06),
                     fill=color)
            add_text(s, x + Inches(0.35), iy, col_w - Inches(0.5), Inches(0.4),
                     it, size=10.5, color=BRAND["grey_dark"], line_spacing=1.4)

    # Handoff strip
    hy = Inches(6.45)
    add_rect(s, Inches(0.5), hy, Inches(12.33), Inches(0.7), fill=BRAND["lightgrey"])
    add_rect(s, Inches(0.5), hy, Inches(0.1), Inches(0.7), fill=BRAND["red"])
    add_text(s, Inches(0.8), hy + Inches(0.12), Inches(12), Inches(0.22),
             "CRITICAL HANDOFFS", size=9, bold=True, color=BRAND["red"], tracking=160)
    add_text(s, Inches(0.8), hy + Inches(0.34), Inches(12), Inches(0.3),
             "RA → procurement (Day 30) · Site ready → racking (Day 60) · "
             "Beta accepted → GA cut-over (Day 90) · GA → managed service (Month 4)",
             size=10.5, color=BRAND["grey_dark"])

    add_footer(s, 8)


def slide_risks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        s, "Slide 09 · Risks, assumptions, decisions",
        "What we need from the room",
        ("Five risks tracked, four assumptions to validate, three decisions "
         "to approve so Fikreab can finalise the commercial deck."))

    # Risks (left half)
    rx = Inches(0.5); ry = Inches(2.55); rw = Inches(6.2)
    add_text(s, rx, ry, rw, Inches(0.3),
             "TOP RISKS & MITIGATIONS", size=10, bold=True, color=BRAND["red"],
             upper=True, tracking=140)
    risks = [
        ("Boundary drift between AIdeology platform ops & Core42 cloud ops",
         "Single-owner rule per overlap zone + named escalation path."),
        ("Customer SLA gaps between platform and infra",
         "e& fronts the customer SLA; back-to-back SLAs to AIdeology & Core42."),
        ("Hardware lead times slip GA date",
         "Pre-vetted NVIDIA-Certified BoMs + phased pod delivery (Q1 → Q3)."),
        ("Pricing overlap on managed services",
         "Three distinct SKUs (platform / infra / customer-facing) with no overlap."),
        ("Security / sovereignty audit blockers",
         "Help AG wrap, NESA evidence packs, T3/T4 trust tiers pre-built."),
    ]
    for j, (r, m) in enumerate(risks):
        iy = ry + Inches(0.4) + Inches(0.78) * j
        add_rect(s, rx, iy, rw, Inches(0.7), fill=BRAND["white"], line=BRAND["border"])
        add_rect(s, rx, iy, Inches(0.08), Inches(0.7), fill=BRAND["red"])
        add_text(s, rx + Inches(0.2), iy + Inches(0.08), rw - Inches(0.3), Inches(0.28),
                 r, size=10.5, bold=True, color=BRAND["black"], line_spacing=1.3)
        add_text(s, rx + Inches(0.2), iy + Inches(0.38), rw - Inches(0.3), Inches(0.3),
                 f"→ {m}", size=9.5, color=BRAND["grey_mid"], line_spacing=1.4)

    # Assumptions + Decisions (right half)
    ax = Inches(6.95); ay = Inches(2.55); aw = Inches(5.88)
    add_text(s, ax, ay, aw, Inches(0.3),
             "KEY ASSUMPTIONS", size=10, bold=True, color=BRAND["red"],
             upper=True, tracking=140)
    assumptions = [
        "e& is the contracting customer & commercial front.",
        "Core42 provides sovereign GPU capacity where required.",
        "AIdeology delivers the AI / platform plane end-to-end.",
        "Pricing ranges firm in 21 days post-sign-off.",
    ]
    for j, a in enumerate(assumptions):
        iy = ay + Inches(0.38) + Inches(0.34) * j
        add_rect(s, ax, iy + Inches(0.12), Inches(0.06), Inches(0.06),
                 fill=BRAND["red"])
        add_text(s, ax + Inches(0.18), iy, aw - Inches(0.2), Inches(0.32),
                 a, size=11, color=BRAND["grey_dark"], line_spacing=1.4)

    # Decisions block
    dy = Inches(4.55)
    add_text(s, ax, dy, aw, Inches(0.3),
             "DECISIONS NEEDED", size=10, bold=True, color=BRAND["red"],
             upper=True, tracking=140)
    decisions = [
        ("01", "Approve RACI"),
        ("02", "Approve commercial frame"),
        ("03", "Approve overlap rules"),
    ]
    for j, (n, d) in enumerate(decisions):
        iy = dy + Inches(0.4) + Inches(0.7) * j
        add_rect(s, ax, iy, aw, Inches(0.6), fill=BRAND["black"])
        add_rect(s, ax, iy, Inches(0.5), Inches(0.6), fill=BRAND["red"])
        add_text(s, ax, iy, Inches(0.5), Inches(0.6), n,
                 size=18, bold=True, color=BRAND["white"],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT)
        add_text(s, ax + Inches(0.7), iy + Inches(0.05), aw - Inches(0.8), Inches(0.5),
                 d, size=13, bold=True, color=BRAND["white"],
                 anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 9)


# ─────────────────────────────────────────────────────────────────────────────
# Build presentation

def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)         # 1
    slide_objective(prs)     # 2
    slide_scope_map(prs)     # 3
    slide_raci(prs)          # 4
    slide_overlap(prs)       # 5
    slide_catalogue(prs)     # 6
    slide_commercial(prs)    # 7
    slide_timeline(prs)      # 8
    slide_risks(prs)         # 9

    out = "AIdeology_GPUaaS_Professional_Services.pptx"
    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    build()
